import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime


DARK = RGBColor(0x0a, 0x0a, 0x0f)
ACCENT = RGBColor(0x7c, 0x6a, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)
MUTED = RGBColor(0xc0, 0xc0, 0xd0)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def accent_bar(slide, top=Inches(1.1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), top, W, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_bullet_text(slide, items, left, top, width, height, size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = Pt(size)
        run.font.color.rgb = color


def generate_pptx(analysis):
    idea = analysis.idea
    prs = new_prs()

    # Slide 1: Title
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s, top=Inches(3.5))
    add_text(s, idea.title or 'Startup Idea', Inches(1), Inches(1.5), Inches(11.33), Inches(1.5),
             size=40, bold=True, align=PP_ALIGN.CENTER)
    pitch_short = (idea.pitch_text or '')[:120] + ('…' if len(idea.pitch_text or '') > 120 else '')
    add_text(s, pitch_short, Inches(2), Inches(3.8), Inches(9.33), Inches(1.5),
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, datetime.utcnow().strftime('%B %Y'), Inches(5), Inches(5.5), Inches(3.33), Inches(0.5),
             size=12, color=RGBColor(0x6a, 0x6a, 0x8a), align=PP_ALIGN.CENTER)

    # Slide 2: The Problem
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'The problem', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    items = []
    for p in (analysis.personas or []):
        if p.pain_point:
            items.append(f'{p.name} ({p.role}): {p.pain_point}')
    if not items:
        items = ['No persona data available']
    add_bullet_text(s, items, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5), size=18)

    # Slide 3: The Solution
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'The solution', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    pitch_lines = (idea.pitch_text or '').split('. ')[:3]
    add_bullet_text(s, pitch_lines, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5), size=18)

    # Slide 4: Market Opportunity
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Market opportunity', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    market = next((sc for sc in (analysis.scores or []) if sc.dimension == 'market_size'), None)
    if market:
        add_text(s, f'Market size score: {market.score:.1f}/10', Inches(0.8), Inches(1.4),
                 Inches(11.5), Inches(0.6), size=22, bold=True)
        add_text(s, market.justification or '', Inches(0.8), Inches(2.2),
                 Inches(11.5), Inches(3), size=16, color=MUTED)

    # Slide 5: Business Model
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Business model', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    add_text(s, f'Model: {(idea.business_model or "").title()}', Inches(0.8), Inches(1.4),
             Inches(11.5), Inches(0.6), size=22, bold=True)
    rev = next((sc for sc in (analysis.scores or []) if sc.dimension == 'revenue_potential'), None)
    if rev:
        add_text(s, rev.justification or '', Inches(0.8), Inches(2.2),
                 Inches(11.5), Inches(3), size=16, color=MUTED)

    # Slide 6: Competitive Landscape
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Competitive landscape', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    lines = []
    for c in (analysis.competitors or []):
        lines.append(f'{c.name} ({c.type}, {c.scope}) — {c.hottest_product or ""}')
    if not lines:
        lines = ['No competitor data available']
    add_bullet_text(s, lines, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5), size=15)

    # Slide 7: Differentiation
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Our differentiation', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    diff = next((sc for sc in (analysis.scores or []) if sc.dimension == 'differentiation'), None)
    strengths = list((analysis.swot.strengths if analysis.swot else []) or [])
    items = []
    if diff:
        items.append(f'Differentiation score: {diff.score:.1f}/10 — {diff.justification or ""}')
    items.extend(strengths[:3])
    add_bullet_text(s, items, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5), size=16)

    # Slide 8: Target Customer
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Target customers', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    for i, p in enumerate((analysis.personas or [])[:2]):
        left = Inches(0.5 + i * 6.5)
        add_text(s, f'{p.name}, {p.age}', left, Inches(1.4), Inches(6), Inches(0.5),
                 size=18, bold=True)
        add_text(s, p.role or '', left, Inches(1.9), Inches(6), Inches(0.4),
                 size=14, color=MUTED)
        add_text(s, f'Pain: {p.pain_point or ""}', left, Inches(2.4), Inches(6), Inches(1.2),
                 size=13, color=WHITE)
        add_text(s, f'WTP: {p.willingness_to_pay or ""}', left, Inches(3.8), Inches(6), Inches(0.5),
                 size=13, color=ACCENT)

    # Slide 9: Risks & Mitigations
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s)
    add_text(s, 'Risks & mitigations', Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT)
    threats = list((analysis.swot.threats if analysis.swot else []) or [])
    pivots = list(analysis.pivot_suggestions or [])
    items = [f'Risk: {t}' for t in threats] + [f'Pivot: {p}' for p in pivots]
    if not items:
        items = ['No risk data available']
    add_bullet_text(s, items, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5), size=15)

    # Slide 10: Ask / Next Steps
    s = blank_slide(prs)
    set_bg(s)
    accent_bar(s, top=Inches(3.5))
    add_text(s, 'Ask & next steps', Inches(0.5), Inches(1.0), Inches(12), Inches(1),
             size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, '[Insert funding ask here]', Inches(2), Inches(2.5), Inches(9.33), Inches(0.7),
             size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullet_text(s, [
        'Validate with 10 target customers',
        'Build MVP in 8 weeks',
        'Apply to accelerator / seek seed funding',
    ], Inches(3.5), Inches(3.5), Inches(6.33), Inches(3), size=15)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
